from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_content_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

def add_image_slide(prs, title, image_path):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), height=Inches(5.5))

def create_unsupervised_learning_diagram(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "教師なし学習による異常検知"

    # データ入力の図形
    input_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(2), Inches(1))
    input_shape.text = "X線Raw Data"
    input_shape.fill.solid()
    input_shape.fill.fore_color.rgb = RGBColor(173, 216, 230)

    # 処理の図形
    process_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(2), Inches(2), Inches(1))
    process_shape.text = "教師なし学習\nアルゴリズム"
    process_shape.fill.solid()
    process_shape.fill.fore_color.rgb = RGBColor(144, 238, 144)

    # 出力の図形
    output_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(2), Inches(2), Inches(1))
    output_shape.text = "異常検知結果"
    output_shape.fill.solid()
    output_shape.fill.fore_color.rgb = RGBColor(255, 182, 193)

    # 矢印1
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.ARROW, Inches(3), Inches(2.25), Inches(1), Inches(0.5))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = RGBColor(0, 0, 0)

    # 矢印2
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.ARROW, Inches(6), Inches(2.25), Inches(1), Inches(0.5))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = RGBColor(0, 0, 0)

def main():
    prs = Presentation()

    # タイトルスライド
    add_title_slide(prs, "X線Raw Dataを用いた微小病変検出プロジェクト", "AI開発部 新規プロジェクト提案")

    # メリットスライド
    merits_content = """
    • 画像再構成前のデータを使用するため、情報損失が最小限
    • 微小な異常をより高感度に検出できる可能性
    • 処理時間の短縮（画像再構成プロセスをスキップ）
    • 新たな特徴量の発見につながる可能性
    """
    add_content_slide(prs, "X線Raw Dataを直接学習することのメリット", merits_content)

    # 必要な要素スライド
    requirements_content = """
    1. Raw Dataの効率的な読み込みシステム
    2. ノイズ処理アルゴリズム
    3. 大容量データの高速処理能力
    4. 特徴抽出のための新たな手法
    5. 教師なし学習アルゴリズムの最適化
    """
    add_content_slide(prs, "プロジェクト実現に必要な要素", requirements_content)

    # 協力要請スライド
    cooperation_content = """
    • Raw Dataの構造と読み込み方法に関する知見
    • 効果的なノイズ処理技術の共有
    • 大規模データ処理のためのインフラ設計支援
    • CT画像再構成アルゴリズムの知見を活かした特徴抽出手法の開発
    • プロジェクトの評価指標設定と検証方法の助言
    """
    add_content_slide(prs, "CT開発部への協力要請事項", cooperation_content)

    # 教師なし学習による異常検知のポンチ絵
    create_unsupervised_learning_diagram(prs)

    # プレゼンテーションの保存
    prs.save('xray_project_proposal.pptx')

if __name__ == "__main__":
    main()
